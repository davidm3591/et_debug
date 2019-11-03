from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def brk_if():
    test_stop = input("\nContinue loop (y/n)? ")
    test_stop = test_stop.lower()
    if test_stop == 'y':
        pass
    else:
        sys.exit(None)


pres = Presentation('3316-09-04_passes_stripped.pptx')
slide = pres.slides[0]

for shape in slide.shapes:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        # continue
        print(shape.shape_type)
        brk_if()
    picture = shape
    print(picture._pic.nvPicPr.cNvPr.get('descr'))


# pres = Presentation('3316-09-04_passes_stripped.pptx')
# slide = pres.slides[1]
# shape = slide.shapes[1]
# image = shape.image
# blob = image.blob
# ext = image.ext
# with open(f'image.{ext}', 'wb') as file:
#     file.write(blob)


# pres = Presentation('3316-09-04_passes_stripped.pptx')
# slide = pres.slides[1]
# shape = slide.shapes[1]
# image = shape.image
# blob = image.blob
# ext = image.ext
# with open(f'image.{ext}', 'wb') as file:
#     file.write(blob)

#
# Resources
#
# https://duckduckgo.com/
# ?q=how+do+I+get+slide+image+content+python-pptx&t=ffnt&atb=v188-1&ia=web&iax=qa
#
# https://stackoverflow.com/questions/46463861/get-image-file-names-with-python-pptx
#
# https://stackoverflow.com/questions/56388068/get-a-picture-python-pptx
#
# https://python-pptx.readthedocs.io/en/latest/api/image.html#pptx.parts.image.Image
#
# https://python-pptx.readthedocs.io/en/latest/api/shapes.html#pptx.shapes.picture.Picture.image
#
#




