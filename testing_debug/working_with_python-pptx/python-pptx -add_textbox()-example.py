from pptx import Presentation

#
# Code diagram follows code here
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Build the Title Slide
# 1. Build title slide with slide_layouts
# 2. Add 'title' slide.shape for text
# 3. Add  'subtitle' slide.placeholders for text
# 4. Add text to 'title' slide.shapes as text (title.text)
#    Add text to 'subtitle' slide.placeholders as text (subtitle.text)
#
# Add a Bullet slide
# 5. Build bullet slide with slide.layouts
# 6. Add bullet slide to the presentation with slides.add_slide() method
# 7. Add 'title_shape' with shapes.title for text
# 8. Add 'body_shape' with placeholders[1]
# 9. Add 'title_shape' text
# 10. Add text_frame to body_shape
# 11. Add text to text_frame
# 12. Add next bullet with 'add_paragraph() to text_frame'

# Build the Title Slide with Title and a Subtitle
# 1.
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
# 2.
title = slide.shapes.title
# 3.
subtitle = slide.placeholders[1]
# 4.
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"


# Build a Bullet Slide
# 5.
bullet_slide_layout = prs.slide_layouts[1]
# 6.
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
# 7.
title_shape = shapes.title
# 8.
body_shape = shapes.placeholders[1]
# 9.
title_shape.text = 'Adding a Bullet Slide'
# 10.
tf = body_shape.text_frame
# 11.
tf.text = 'Find the bullet slide layout'

# 12.
p = tf.add_paragraph()
p.text = 'Use TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2


# Add a textbox() with add_textbox()



prs.save('test.pptx')


#
# Code explanation
#

title_diagram = """
--------------------------------------------------------------
BUILD A NEW POWERPOINT PRESENTATION (CONTINUTED)

ADD A TEXTBOX EXAMPLE:
    slide.shapes.add_textbox(left, top, width, height)
--------------------------------------------------------------

CODE:
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

CODE DIAGRAMMED
 ________________________________________
|    ________________________________    |
|   |                                |   |<--- prs = Presentation()
|   |                                |   |
|   |       ---------------          |<--|---- title_slide_layout = prs.slide_layouts[0]
|   |       | Hello World | <--------|---|---- title = slide.shapes.title
|   |       --------------- <--------|---|---- title.text = 'Hello World'
|   |                                |   |
|   |                                |   |
|   |   ------------------------- <--|---|---- subtitle = slide.placeholders[1]
|   |   | python-pptx was here! | <--|---|---- subtitle.text = 'python-pptx was here!'
|   |   -------------------------    |   |
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |________________________________|   |
|________________________________________|

"""

# print(title_diagram)


bullet_diagram = # --------------------------------------------------------------
# BUILD A NEW POWERPOINT PRESENTATION (CONTINUED)

# BUILD A NEW BULLET POINT SLIDE:
#    A TILE SLIDE AND TEXT (SHAPE)
#    A SUBTITLE PLACEHOLER WITH TEXT
# --------------------------------------------------------------

# CODE:
# prs = Presentation()

# bullet_slide_layout = prs.slide_layouts[1]
# slide = prs.slides.add_slide(bullet_slide_layout)
# shapes = slide.shapes

# title_shape = shapes.title
# body_shape = shapes.placeholders[1]

# title_shape.text = 'Adding a Bullet Slide'

# tf = body_shape.text_frame
# tf.text = 'Find the bullet slide layout'

# p = tf.add_paragraph()
# p.text = 'Use TextFrame.text for first bullet'
# p.level = 1

# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
# p.level = 2

# prs.save('test.pptx')

# CODE DIAGRAMMED
#  _________________________________________________________
# |    ________________________________________________     |
# |   |                                                |    |<--- prs = Presentation()
# |   |                                                |    |
# |   |   ------------------------- <------------------|----|---- bullet_slide_layout
# |   |   | Adding a Bullet Slide |                    |    |      = prs.slide_layouts[1]
# |   |   |                       | <------------------|----|---- title_shape = shapes.title
# |   |   ------------------------- <------------------|----|---- title_shape.text = 'Hello World'
# |   |                                                |    |
# |   |   ------------------------- <------------------|----|---- body_shape = shapes.placeholders[1]
# |   |   | - Find bullet slide   | <------------------|----|---- tf = body_shape.text_frame
# |   |   ------------------------- <------------------|----|---- tf.text = 'bullet slide layout'
# |   |                                                |    |
# |   |     ------------------------ <-----------------|----|---- p = tf.add_paragraph()
# |   |     | x Use TextFrame.text | <-----------------|----|---- p.text = 'Use TextFrame.text'
# |   |     |   for first bullet   |                   |    |       for first bullet
# |   |     |                      | <-----------------|----|---- p.level = 1
# |   |     ------------------------                   |    |
# |   |                                                |    |
# |   |       ----------------------------------- <----|----|---- p = tf.add_paragraph()
# |   |       | * Use TextFrame.add_paragraph() | <----|----|---- p.text = 'Use _TextFrame.add_paragraph()
# |   |       |   for subsequent bullets        |      |    |       for subsequent bullets'
# |   |       |                                 | <----|----|---- p.level = 2
# |   |       -----------------------------------      |    |
# |   |                                                |    |
# |   |________________________________________________|    |
# |_________________________________________________________|

#

# print(title_diagram, bullet_diagram)
