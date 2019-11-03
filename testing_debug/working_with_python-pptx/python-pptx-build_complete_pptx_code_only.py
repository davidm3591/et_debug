from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
#
# Code diagram follows code here
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Build the Title Slide
# 1. Build title slide with slide_layouts
# 2. Add 'title' slide.shape for text
# 3. Add  'subtitle' slide.placeholders for text
# 4. Add text to 'title' slide.shapes as text (title.text)
#    Add text to 'subtitle' slide.placeholders as text (subtitle.text)
#
# Add a Bullet slide
# 5. Build bullet slide with slide.layouts
# 6. Add bullet slide to the presentation with slides.add_slide() method
# 7. Add 'title_shape' with shapes.title for text
# 8. Add 'body_shape' with placeholders[1]
# 9. Add 'title_shape' text
# 10. Add text_frame to body_shape
# 11. Add text to text_frame
# 12. Add next bullet with 'add_paragraph() to text_frame'
#
# Add a textbox() with add_textbox()
# 13. Add blank slide with slide.layouts[6]
# 14. Add slide for textbox with slides.add_slide() method
# 15. Dimension textbox
# 16. Add the textbox with slide.shapes.add_textbox() method
# 17. Add text to textbox shape
# 18. Add more text in another paragraph with add_paragraph() method
#     Bold font format with font.bold = True
# 19. Add more text in another paragraph with add_paragraph() method
#     Large font size with font.size(40)
#
# Add a picture with add_picture()
# 20. Create image path
# 21. Add blank slide with slide.layouts[6]
# 22. Add slide for image/picture with slides.add_slide() method
# 23. Dimension slide shape
# 24. Add picture
# 25. Add, position, size picture
#
#
# Add an AutoShape
# 26. Create the slide layout
# 27. Add slide to presentation
# 28. Add a shape for the Title text
# 29. Add text to title shape
# 30. Position and size the AutoShape
# 31. Add shape and shape text
# 32. Resize shape for visual balance
#     Make chevrons wider
# 33. Add for loop to build chevrons
# 
# 
# Add a Table
# 34. New slide
# 35. Add title text
# 36. Define table, cols, size, position
# 37. Create table object
# 38. Set column widths
# 39. Write column headings
# 40. Write body cells
#################################################


#
# Build the Title Slide with Title and a Subtitle
#
# 1.
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
# 2.
title = slide.shapes.title
# 3.
subtitle = slide.placeholders[1]
# 4.
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"


#
# Build a Bullet Slide
#
# 5.
bullet_slide_layout = prs.slide_layouts[1]
# 6.
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
# 7.
title_shape = shapes.title
# 8.
body_shape = shapes.placeholders[1]
# 9.
title_shape.text = 'Adding a Bullet Slide'
# 10.
tf = body_shape.text_frame
# 11.
tf.text = 'Find the bullet slide layout'

# 12.
p = tf.add_paragraph()
p.text = 'Use TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2


#
# Add a textbox() with add_textbox()
#
from pptx.util import Inches, Pt
# 13.
blank_slide_layout = prs.slide_layouts[6]
# 14.
slide = prs.slides.add_slide(blank_slide_layout)
# 15.
left = top = width = height = Inches(1)
# 16.
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
# 17.
tf.text = "This is text inside a textbox"

# 18.
p = tf.add_paragraph()
p.text = "This is a second paragraph that is bold"
p.font.bold = True

# 19.
p = tf.add_paragraph()
p.text = "This is a third paragraph that is large font"
p.font.size = Pt(40)


#
# Add a picture with add_picture()
#
# 20.
img_path = 'add-picture.png'
# 21.
blank_slide_layout = prs.slide_layouts[6]
# 22.
slide = prs.slides.add_slide(blank_slide_layout)
# 23.
left = top = Inches(1)
# 24.
pic = slide.shapes.add_picture(img_path, left, top)
#
# Add a second picture to same slide
#
# 25.
left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)


#
# Add an AutoShape
#
# 26.
title_only_slide_layout = prs.slide_layouts[5]
# 27.
slide = prs.slides.add_slide(title_only_slide_layout)
# 28.
shapes = slide.shapes
# 29.
shapes.title.text = 'Adding an AutoShape'

# 30.
left = Inches(0.93)  # centers overall set
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

# 31. Add shape and shape text
shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = 'Step 1'

# 32.
left = left + width - Inches(0.4)
width = Inches(2.0)  # Make chevrons wider

# 33.
for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = f'Step {n}'
    left = left + width - Inches(0.4)

# Constants representing each of the available auto shapes
# (like MSO_SHAPE.ROUNDED_RECT, MSO_SHAPE.CHEVRON, etc.)
# are listed on the autoshape - types
# (https://python-pptx.readthedocs.io/en/latest/api/enum
#     /MsoAutoShapeType.html\#msoautoshapetype) page.

#
# Add a Table
#
# 34. New slide
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

# 35. Add title text
shapes.title.text = 'Adding a Table'

# 36. Define table, cols, size, position
rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

# 37. Create table object
table = shapes.add_table(rows, cols, left, top, width, height).table

# 38. Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# 39. Write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# 40. Write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

prs.save('test.pptx')
