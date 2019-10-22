from pptx import Presentation

#
# Code diagram follows code here
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Build the Title Slide
# 1. Build title slide with slide_layouts
# 2. Add 'title' slide.shape for text
# 3. Add  'subtitle' slide.placeholders for text
# 4. Add text to 'title' slide.shapes as text (title.text)
#    Add text to 'subtitle' slide.placeholders as text (subtitle.text)
#
# Add a Bullet slide
# 5. Build bullet slide with slide.layouts
# 6. Add bullet slide to the presentation with slides.add_slide() method

# Build the Title Slide with Title and a Subtitle
# 1.
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
# 2.
title = slide.shapes.title
# 3.
subtitle = slide.placeholders[1]
# 4.
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"


# Build a Bullet Slide
# 5.
bullet_slide_layout = prs.slide_layouts[1]
# 6.
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2


prs.save('test.pptx')


#
# Code explanation
#

title_diagram = """
--------------------------------------------------------------
BUILD A NEW POWERPOINT PRESENTATION

BUILD A TITLE SLIDE
   A TILE SLIDE AND TEXT (SHAPE)
   A SUBTITLE PLACEHOLER WITH TEXT
--------------------------------------------------------------

CODE:
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

CODE DIAGRAMMED
 ________________________________________
|    ________________________________    |
|   |                                |   |<--- prs = Presentation()
|   |                                |   |
|   |       ---------------          |<--|---- title_slide_layout = prs.slide_layouts[0]
|   |       | Hello World | <--------|---|---- title = slide.shapes.title
|   |       --------------- <--------|---|---- title.text = 'Hello World'
|   |                                |   |
|   |                                |   |
|   |   ------------------------- <--|---|---- subtitle = slide.placeholders[1]
|   |   | python-pptx was here! | <--|---|---- subtitle.text = 'python-pptx was here!'
|   |   -------------------------    |   |
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |________________________________|   |
|________________________________________|

"""

print(title_diagram)


bullet_diagram = """
--------------------------------------------------------------
BUILD A NEW POWERPOINT PRESENTATION (CONTINUED)

BUILD A NEW BULLET POINT SLIDE:
   A TILE SLIDE AND TEXT (SHAPE)
   A SUBTITLE PLACEHOLER WITH TEXT
--------------------------------------------------------------

CODE:
prs = Presentation()

bullet_slide_layout = prs.slide_layouts[1]
# 6.
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')

CODE DIAGRAMMED
 ________________________________________
|    ________________________________    |
|   |                                |   |<--- prs = Presentation()
|   |                                |   |
|   |   -------------------------    |<--|---- bullet_slide_layout = prs.slide_layouts[1]
|   |   | Adding a Bullet Slide | <--|---|---- title_shape = shapes.title
|   |   ------------------------- <--|---|---- title_shape.text = 'Hello World'
|   |                                |   |
|   |   ------------------------- <--|---|---- body_shape = shapes.placeholders[1]
|   |   | X Find bullet slide   | <--|---|---- tf = body_shape.text_frame
|   |   ------------------------- <--|---|---- tf.text = 'bullet slide layout'
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |                                |   |
|   |________________________________|   |
|________________________________________|

"""

print(bullet_diagram)
