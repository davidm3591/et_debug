from pptx import Presentation

#
# Build a Title Slide with Title and Subtitle
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Build the Title Slide
#   1. Build title slide with slide_layouts
#   2. Add 'title' slide.shape for text
#   3. Add  'subtitle' slide.placeholders for text
#   4. Add text to 'title' slide.shapes as text (title.text)
#      Add text to 'subtitle' slide.placeholders as text (subtitle.text)
#

#
# Build the Title Slide with Title and a Subtitle
#
# 1.
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
# 2.
title = slide.shapes.title
# 3.
subtitle = slide.placeholders[1]
# 4.
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"


prs.save('test.pptx')
