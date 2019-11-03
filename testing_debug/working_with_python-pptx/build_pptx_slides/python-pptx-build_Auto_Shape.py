from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


#
# Build a AutoShape Slide
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Add an AutoShape
#   26. Create the slide layout
#   27. Add slide to presentation
#   28. Add a shape for the Title text
#   29. Add text to title shape
#   30. Position and size the AutoShape
#   31. Add shape and shape text
#   32. Resize shape for visual balance
#       Make chevrons wider
#   33. Add for loop to build chevrons
#


#
# Add an AutoShape
#
# 26.
title_only_slide_layout = prs.slide_layouts[5]
# 27.
slide = prs.slides.add_slide(title_only_slide_layout)
# 28.
shapes = slide.shapes
# 29.
shapes.title.text = 'Adding an AutoShape'

# 30.
left = Inches(0.93)  # centers overall set
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

# 31. Add shape and shape text
shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = 'Step 1'

# 32.
left = left + width - Inches(0.4)
width = Inches(2.0)  # Make chevrons wider

# 33.
for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = f'Step {n}'
    left = left + width - Inches(0.4)

# Constants representing each of the available auto shapes
# (like MSO_SHAPE.ROUNDED_RECT, MSO_SHAPE.CHEVRON, etc.)
# are listed on the autoshape - types
# (https://python-pptx.readthedocs.io/en/latest/api/enum
#     /MsoAutoShapeType.html\#msoautoshapetype) page.


prs.save('test.pptx')
