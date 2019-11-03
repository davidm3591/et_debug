from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


#
# Build a Bullet Slide with Different Levels
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Build a Bullet slide
#   5. Build bullet slide with slide.layouts
#   6. Add bullet slide to the presentation with slides.add_slide() method
#   7. Add 'title_shape' with shapes.title for text
#   8. Add 'body_shape' with placeholders[1]
#   9. Add 'title_shape' text
#   10. Add text_frame to body_shape
#   11. Add text to text_frame
#   12. Add next bullet with 'add_paragraph() to text_frame'
#  

#
# Build a Bullet Slide and multiple list levels
#
# 5.
bullet_slide_layout = prs.slide_layouts[1]
# 6.
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
# 7.
title_shape = shapes.title
# 8.
body_shape = shapes.placeholders[1]
# 9.
title_shape.text = 'Adding a Bullet Slide'
# 10.
tf = body_shape.text_frame
# 11.
tf.text = 'Find the bullet slide layout'

# 12.
p = tf.add_paragraph()
p.text = 'Use TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2


prs.save('test.pptx')
