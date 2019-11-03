from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


#
# Build an Image Slide
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Add a picture with add_picture()
#   20. Create image path
#   21. Add blank slide with slide.layouts[6]
#   22. Add slide for image/picture with slides.add_slide() method
#   23. Dimension slide shape
#   24. Add picture
#   25. Add, position, size picture
#


#
# Add a picture with add_picture()
#
# 20.
img_path = 'add-picture.png'
# 21.
blank_slide_layout = prs.slide_layouts[6]
# 22.
slide = prs.slides.add_slide(blank_slide_layout)
# 23.
left = top = Inches(1)
# 24.
pic = slide.shapes.add_picture(img_path, left, top)
#
# Add a second picture to same slide
#
# 25.
left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
