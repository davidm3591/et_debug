from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


#
# Build a Table Slide
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Add a Table
#   34. New slide
#   35. Add title text
#   36. Define table, cols, size, position
#   37. Create table object
#   38. Set column widths
#   39. Write column headings
#   40. Write body cells
#


#
# Add a Table
#
# 34. New slide
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

# 35. Add title text
shapes.title.text = 'Adding a Table'

# 36. Define table, cols, size, position
rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

# 37. Create table object
table = shapes.add_table(rows, cols, left, top, width, height).table

# 38. Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# 39. Write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# 40. Write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

prs.save('test.pptx')
