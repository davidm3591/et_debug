from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


#
# Build a Textbox Slide
#

prs = Presentation()

# ####################################################################
#
# Building a PowerPoint
#
# Add a textbox() with add_textbox()
#   13. Add blank slide with slide.layouts[6]
#   14. Add slide for textbox with slides.add_slide() method
#   15. Dimension textbox
#   16. Add the textbox with slide.shapes.add_textbox() method
#   17. Add text to textbox shape
#   18. Add more text in another paragraph with add_paragraph() method
#       Bold font format with font.bold = True
#   19. Add more text in another paragraph with add_paragraph() method
#       Large font size with font.size(40)
#


#
# Add a textbox() with add_textbox()
#
from pptx.util import Inches, Pt

prs = Presentation()

# 13.
blank_slide_layout = prs.slide_layouts[6]
# 14.
slide = prs.slides.add_slide(blank_slide_layout)
# 15.
left = top = width = height = Inches(1)
# 16.
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
# 17.
tf.text = "This is text inside a textbox"

# 18.
p = tf.add_paragraph()
p.text = "This is a second paragraph that is bold"
p.font.bold = True

# 19.
p = tf.add_paragraph()
p.text = "This is a third paragraph that is large font"
p.font.size = Pt(40)

prs.save('test.pptx')